"""Product Overview Deck — assembled from real analysis data via python-pptx.

Six slides, brand-styled:
    1. Cover — site name, score banner, verdict, key facts
    2. Score Glass-Box breakdown chart + dimension table
    3. Proximity intelligence — static map + category counts chart
    4. Risk overlay — flood-risk callout + sustainability table
    5. Pro-Forma — cashflow chart + Bear/Base/Bull comparison
    6. Methodology + sources

Charts are PNG bytes from `utils.charts`. No browser automation; the deck
reflects whatever analysis is in `site`.
"""
from __future__ import annotations

from datetime import datetime
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt, Emu

from utils import charts


# ---------------------------------------------------------------------------
# Style constants
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)
ACCENT_RGB = RGBColor(0x00, 0xC8, 0x53)
PRIMARY_RGB = RGBColor(0x1A, 0x1A, 0x1A)
SUBTLE_RGB = RGBColor(0x55, 0x55, 0x55)
LIGHT_RGB = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]   # Blank
    return prs.slides.add_slide(layout)


def _rect(slide, left, top, width, height, fill_rgb, line=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if not line:
        shape.line.fill.background()
    return shape


def _text(slide, left, top, width, height, text, *,
           size=14, bold=False, color=PRIMARY_RGB, align="left"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": 1, "center": 2, "right": 3}.get(align, 1)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_image(slide, png_bytes: bytes, left, top, *, width=None, height=None):
    return slide.shapes.add_picture(
        BytesIO(png_bytes), left, top, width=width, height=height,
    )


def _slide_header(slide, brand: dict, title: str, page_num: int, total: int):
    # Brand band top
    _rect(slide, 0, 0, SLIDE_W, Inches(0.55), PRIMARY_RGB)
    _text(slide, Inches(0.4), Inches(0.08), Inches(8), Inches(0.4),
           f"{brand['brand']}  ·  Retail Site Intelligence  ·  Product Overview",
           size=11, color=WHITE)
    _text(slide, Inches(11.0), Inches(0.08), Inches(2), Inches(0.4),
           f"{page_num} / {total}", size=11, color=WHITE, align="right")
    # Page title
    _text(slide, Inches(0.4), Inches(0.7), Inches(12.5), Inches(0.6),
           title, size=22, bold=True, color=PRIMARY_RGB)
    # Accent underline
    _rect(slide, Inches(0.4), Inches(1.35), Inches(0.6), Inches(0.05),
           ACCENT_RGB)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def _verdict_for(esi: int) -> tuple[str, RGBColor]:
    if esi >= 75:
        return "🟢 Go", RGBColor(0x00, 0xC8, 0x53)
    if esi >= 50:
        return "🟡 Vertiefte Prüfung", RGBColor(0xF9, 0xA8, 0x25)
    return "🔴 No-Go", RGBColor(0xD3, 0x2F, 0x2F)


def _slide_cover(prs, brand, site, total):
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.55), PRIMARY_RGB)
    _rect(slide, 0, Inches(7.0), SLIDE_W, Inches(0.5), PRIMARY_RGB)
    _text(slide, Inches(0.4), Inches(0.08), Inches(8), Inches(0.4),
           f"{brand['brand']}  ·  Product Overview Deck", size=12, color=WHITE)
    _text(slide, Inches(11.0), Inches(0.08), Inches(2), Inches(0.4),
           datetime.now().strftime("%Y-%m-%d"),
           size=12, color=WHITE, align="right")

    _text(slide, Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.6),
           "Retail Site Intelligence", size=32, bold=True, color=PRIMARY_RGB)
    _text(slide, Inches(0.4), Inches(1.7), Inches(12.5), Inches(0.5),
           "Site selection · pipeline tracking · portfolio management",
           size=15, color=SUBTLE_RGB)

    _rect(slide, Inches(0.4), Inches(2.5), Inches(0.6), Inches(0.06),
           ACCENT_RGB)
    _text(slide, Inches(0.4), Inches(2.6), Inches(12.5), Inches(0.4),
           f"Case study · {site.get('name', '—')}",
           size=20, bold=True, color=PRIMARY_RGB)
    _text(slide, Inches(0.4), Inches(3.05), Inches(12.5), Inches(0.4),
           site.get("address", ""), size=12, color=SUBTLE_RGB)

    # Score banner
    score = site.get("score_esi") or site.get("score") or 0
    verdict_text, verdict_color = _verdict_for(score)
    _rect(slide, Inches(0.4), Inches(3.7), Inches(3.0), Inches(2.2),
           LIGHT_RGB)
    _text(slide, Inches(0.55), Inches(3.85), Inches(2.8), Inches(0.35),
           "ESI-ADJUSTED SCORE", size=10, bold=True, color=SUBTLE_RGB)
    _text(slide, Inches(0.55), Inches(4.2), Inches(2.8), Inches(1.0),
           str(score), size=64, bold=True, color=verdict_color)
    _text(slide, Inches(0.55), Inches(5.2), Inches(2.8), Inches(0.4),
           verdict_text, size=14, bold=True, color=PRIMARY_RGB)

    # Right-side key facts
    facts_left = Inches(3.8)
    facts_top = Inches(3.7)
    result = site.get("score_result")
    flood = site.get("flood_risk")
    sust = site.get("sustainability_result")
    facts = [
        ("Premium Env. Index", f"{result.pei:.1f} / 10" if result else "—"),
        ("Flood-risk", flood.label if flood else "—"),
        ("Sustainability Δ",
         (f"{sust.score_delta:+d} pts · rent {sust.rent_delta_pct:+.1f}%"
          if sust else "—")),
        ("Pro-Forma NPV",
         f"CHF {int(site['proforma'].npv):,}".replace(",", " ")
         if site.get("proforma") else "—"),
        ("Pro-Forma IRR",
         f"{site['proforma'].irr * 100:.1f}%"
         if site.get("proforma") and site["proforma"].irr else "—"),
    ]
    for i, (k, v) in enumerate(facts):
        _text(slide, facts_left, facts_top + Inches(0.45 * i),
               Inches(3.2), Inches(0.4), k, size=10, color=SUBTLE_RGB)
        _text(slide, facts_left + Inches(3.2), facts_top + Inches(0.45 * i),
               Inches(6.0), Inches(0.4), v, size=13, bold=True,
               color=PRIMARY_RGB)

    _text(slide, Inches(0.4), Inches(7.08), Inches(12.5), Inches(0.4),
           "Methodology: OSM · BFS · BAFU · ESI (CCRS) · Bechtiger 2024",
           size=9, color=WHITE)
    return slide


def _slide_score(prs, brand, site, page_num, total):
    slide = _blank_slide(prs)
    _slide_header(slide, brand, "Site Score · Glass Box breakdown",
                   page_num, total)
    # Chart on left
    result = site.get("score_result")
    if result:
        flood = site.get("flood_risk")
        sust = site.get("sustainability_result")
        png = charts.chart_score_breakdown(
            result,
            flood_penalty=flood.score_impact if flood else 0,
            sust_delta=sust.score_delta if sust else 0,
        )
        _add_image(slide, png, Inches(0.4), Inches(1.6), width=Inches(8.0))

    # Methodology callout on right
    _rect(slide, Inches(8.7), Inches(1.6), Inches(4.3), Inches(5.2),
           LIGHT_RGB)
    _text(slide, Inches(8.9), Inches(1.75), Inches(4.0), Inches(0.4),
           "Methodology", size=14, bold=True, color=PRIMARY_RGB)
    _rect(slide, Inches(8.9), Inches(2.15), Inches(0.4), Inches(0.05),
           ACCENT_RGB)

    text_blocks = [
        ("Every dimension normalised 0–100 BEFORE the weighted sum, so "
         "each component is directly comparable across sites."),
        ("Weights are user-editable via the sidebar; the same engine "
         "produces a defensible score for any retailer (just swap "
         "on_brand.json)."),
        ("Flood penalty is applied POST-aggregation. "
         "Sustainability Δ (Bechtiger 2024) is additive outside the 100-pt "
         "base."),
    ]
    y = Inches(2.3)
    for block in text_blocks:
        _text(slide, Inches(8.9), y, Inches(4.0), Inches(1.4),
               block, size=10, color=PRIMARY_RGB)
        y += Inches(1.35)
    return slide


def _slide_proximity(prs, brand, site, page_num, total):
    slide = _blank_slide(prs)
    _slide_header(slide, brand, "Proximity intelligence · OSM signal",
                   page_num, total)

    pois = site.get("pois") or []
    # Static map on left
    if pois and site.get("lat"):
        try:
            map_png = charts.chart_static_map(
                pois, site["lat"], site["lon"],
                site.get("radius_m", 500),
                title=f"Site context · {site.get('radius_m', 500)} m",
            )
            _add_image(slide, map_png, Inches(0.4), Inches(1.6),
                        width=Inches(6.4))
        except Exception:
            pass

    # POI category chart on right
    if pois:
        try:
            cat_png = charts.chart_poi_categories(
                pois, site.get("radius_m", 500),
            )
            _add_image(slide, cat_png, Inches(7.1), Inches(1.6),
                        width=Inches(5.9))
        except Exception:
            pass

    # Footer narrative
    summary = {}
    for p in pois:
        s = summary.setdefault(p.category, {"count": 0, "score": 0.0})
        s["count"] += 1
        s["score"] += p.score
    parts = [
        f"{cat.title()} {summary[cat]['count']} ({summary[cat]['score']:+.0f})"
        for cat in ("sport", "symbiose", "premium", "competitor", "negative")
        if cat in summary
    ]
    _text(slide, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.5),
           " · ".join(parts), size=10, color=SUBTLE_RGB, align="left")
    return slide


def _slide_risk(prs, brand, site, page_num, total):
    slide = _blank_slide(prs)
    _slide_header(slide, brand, "Risk overlay · Flood + Sustainability",
                   page_num, total)

    # Left panel: flood
    fr = site.get("flood_risk")
    _rect(slide, Inches(0.4), Inches(1.6), Inches(6.2), Inches(5.4),
           LIGHT_RGB)
    _text(slide, Inches(0.6), Inches(1.75), Inches(5.8), Inches(0.4),
           "🌊 BAFU flood-risk", size=14, bold=True, color=PRIMARY_RGB)
    _rect(slide, Inches(0.6), Inches(2.15), Inches(0.4), Inches(0.05),
           ACCENT_RGB)
    if fr:
        _text(slide, Inches(0.6), Inches(2.3), Inches(5.8), Inches(0.6),
               fr.label, size=12, bold=True, color=PRIMARY_RGB)
        _text(slide, Inches(0.6), Inches(3.0), Inches(5.8), Inches(2.6),
               fr.recommendation, size=10, color=PRIMARY_RGB)
        _text(slide, Inches(0.6), Inches(5.8), Inches(5.8), Inches(0.4),
               f"Score impact: {fr.score_impact:+d} pts", size=11,
               bold=True, color=PRIMARY_RGB)
        _text(slide, Inches(0.6), Inches(6.2), Inches(5.8), Inches(0.4),
               (f"Business-interruption insurance estimate: "
                f"CHF {fr.insurance_chf:,}/yr").replace(",", " "),
               size=11, color=PRIMARY_RGB)
        _text(slide, Inches(0.6), Inches(6.55), Inches(5.8), Inches(0.4),
               f"Source: {fr.source}", size=8, color=SUBTLE_RGB)
    else:
        _text(slide, Inches(0.6), Inches(2.3), Inches(5.8), Inches(0.6),
               "Flood check not run for this site.", size=12, color=SUBTLE_RGB)

    # Right panel: sustainability
    sust = site.get("sustainability_result")
    _rect(slide, Inches(6.8), Inches(1.6), Inches(6.2), Inches(5.4),
           LIGHT_RGB)
    _text(slide, Inches(7.0), Inches(1.75), Inches(5.8), Inches(0.4),
           "🌱 Sustainability (Bechtiger 2024 / ESI)",
           size=14, bold=True, color=PRIMARY_RGB)
    _rect(slide, Inches(7.0), Inches(2.15), Inches(0.4), Inches(0.05),
           ACCENT_RGB)
    if sust:
        y = Inches(2.3)
        for f in sust.factors:
            _text(slide, Inches(7.0), y, Inches(2.4), Inches(0.4),
                   f.name, size=11, bold=True, color=PRIMARY_RGB)
            _text(slide, Inches(7.0), y + Inches(0.35),
                   Inches(5.5), Inches(0.4),
                   f.label, size=9, color=SUBTLE_RGB)
            _text(slide, Inches(9.4), y, Inches(1.3), Inches(0.4),
                   f"{f.score_delta:+d} pts", size=11, bold=True,
                   color=PRIMARY_RGB, align="right")
            _text(slide, Inches(10.7), y, Inches(2.0), Inches(0.4),
                   f"{f.rent_delta_pct:+.1f}% rent" if f.rent_delta_pct
                   else "n.s. for rent",
                   size=10, color=SUBTLE_RGB, align="right")
            y += Inches(0.85)
        _text(slide, Inches(7.0), Inches(6.55), Inches(5.8), Inches(0.4),
               f"Total: {sust.score_delta:+d} pts  ·  "
               f"rent {sust.rent_delta_pct:+.1f}%",
               size=11, bold=True, color=PRIMARY_RGB)
    return slide


def _slide_proforma(prs, brand, site, page_num, total):
    slide = _blank_slide(prs)
    _slide_header(slide, brand, "Pro-Forma · 5-year DCF",
                   page_num, total)
    pf = site.get("proforma")
    sc = site.get("scenarios")
    if pf:
        try:
            cf_png = charts.chart_cashflow(pf)
            _add_image(slide, cf_png, Inches(0.4), Inches(1.6),
                        width=Inches(6.4))
        except Exception:
            pass
    if sc:
        try:
            sc_png = charts.chart_scenarios(sc)
            _add_image(slide, sc_png, Inches(7.0), Inches(1.6),
                        width=Inches(6.0))
        except Exception:
            pass

    # Metrics strip at bottom
    if pf:
        fmt = lambda v: f"CHF {int(v):,}".replace(",", " ")
        metrics = [
            ("NPV", fmt(pf.npv)),
            ("IRR",
             f"{pf.irr * 100:.1f}%" if pf.irr is not None else "n/a"),
            ("Payback",
             f"{pf.payback_years:.1f} yr" if pf.payback_years else ">5y"),
            ("Break-even",
             f"M{pf.breakeven_month}" if pf.breakeven_month else "—"),
        ]
        col_w = Inches(2.95)
        for i, (k, v) in enumerate(metrics):
            x = Inches(0.4) + col_w * i + Inches(0.1 * i)
            _rect(slide, x, Inches(6.0), col_w, Inches(1.05), LIGHT_RGB)
            _text(slide, x + Inches(0.15), Inches(6.1),
                   col_w - Inches(0.3), Inches(0.3), k, size=10, bold=True,
                   color=SUBTLE_RGB)
            _text(slide, x + Inches(0.15), Inches(6.4),
                   col_w - Inches(0.3), Inches(0.55), v, size=18, bold=True,
                   color=PRIMARY_RGB)
    return slide


def _slide_methodology(prs, brand, site, page_num, total):
    slide = _blank_slide(prs)
    _slide_header(slide, brand, "Methodology · data sources",
                   page_num, total)
    items = [
        ("OpenStreetMap (OSMnx)",
         "Single query returns sport venues, retail brands, hotels, "
         "transit stops; classification is brand-aware (Lululemon →"
         " symbiose, Nike → competitor)."),
        ("BFS STATPOP + regional Kaufkraft",
         "Population, growth, age structure, purchasing-power index "
         "seeded for 24 ZH-relevant municipalities (extensible CSV)."),
        ("BAFU surface-runoff hazard layer",
         "WMS overlay on the map; analyst-authoritative override for "
         "the four real On CH locations (per ANALYSE file)."),
        ("Pandana-free Walk Score",
         "Amenity richness within walking distance, log-saturated by "
         "category. Pragmatic equivalent of the original Walk Score™ "
         "concept without Windows-hostile dependencies."),
        ("ESI (CCRS/UZH) + Bechtiger 2024",
         "Sustainability factors with empirical coefficients: "
         "accessibility (p<.001), daylight (p=.021), ÖV-class "
         "(positive but n.s. for rent)."),
        ("Glass Box scoring",
         "Every dimension normalised 0–100 before weighting; full "
         "breakdown of what raises/lowers the score is inspectable."),
    ]
    y = Inches(1.6)
    for title, body in items:
        _rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.7),
               ACCENT_RGB)
        _text(slide, Inches(0.7), y, Inches(4.0), Inches(0.4),
               title, size=12, bold=True, color=PRIMARY_RGB)
        _text(slide, Inches(0.7), y + Inches(0.3),
               Inches(12.2), Inches(0.4),
               body, size=9.5, color=SUBTLE_RGB)
        y += Inches(0.85)
    return slide


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def build_overview_deck(brand: dict, site: dict) -> bytes:
    """Return a six-slide product-overview deck as PPTX bytes."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 6
    _slide_cover(prs, brand, site, total)
    _slide_score(prs, brand, site, 2, total)
    _slide_proximity(prs, brand, site, 3, total)
    _slide_risk(prs, brand, site, 4, total)
    _slide_proforma(prs, brand, site, 5, total)
    _slide_methodology(prs, brand, site, 6, total)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
