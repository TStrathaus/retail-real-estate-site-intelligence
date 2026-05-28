"""Product-Overview Deck v2 — embeds real Playwright screenshots.

Companion to `tools/make_screenshots.py`. Drops the matplotlib-chart
slides of v1 and uses the actual rendered UI captures. Run order:

    1. streamlit run app.py            # boots on :8765
    2. python tools/make_screenshots.py  # writes outputs/screenshots/*.png
    3. python -c "from utils.presentation_pptx_v2 import build_screenshot_deck
       from utils.config import load_brand; import pathlib
       pathlib.Path('outputs/ProductOverview_screenshots.pptx').write_bytes(
           build_screenshot_deck(load_brand()))"

Or use sample_run_v2.py (writes both PDFs and PPTs).
"""
from __future__ import annotations

from datetime import datetime
from io import BytesIO
from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ACCENT = RGBColor(0x00, 0xC8, 0x53)
PRIMARY = RGBColor(0x1A, 0x1A, 0x1A)
SUBTLE = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SCREENSHOTS = Path("outputs/screenshots")
LAYER_SHOTS = Path("outputs/proximity_layers")

# Proximity-tab layer screenshots → (file, layer name, description).
# 12 overlays, grouped 4-per-slide. Baseline (no overlay) is omitted —
# it's the reference state, not a layer.
PROXIMITY_LAYERS = [
    ("proximity_01_bafu_flood.jpg", "🌊 BAFU flood hazard",
     "Surface-runoff WMS overlay. Niederdorfstrasse sits in the "
     "Sihl/Limmat HQ100 perimeter — blue tint marks flood-exposed area. "
     "Drives the -15 pt flood penalty + business-interruption insurance."),
    ("proximity_02_running_infrastructure.jpg", "🏃 Running infrastructure",
     "OSM parks (green fill) + running routes / hiking trails / cycleways "
     "(orange lines). Public-data substitute for the Strava heatmap "
     "(~140 features). Limmatufer running corridor visible."),
    ("proximity_03_pedestrian_streets.jpg", "🚶 Pedestrian streets",
     "Fussgängerzonen incl. those where trams + taxis are allowed "
     "(Niederdorf, Limmatquai). Altstadt is an almost continuous "
     "pedestrian ribbon — high spontaneous-footfall context."),
    ("proximity_04_shop_density.jpg", "🔥 Shop-density heatmap",
     "Every OSM shop within 1.5× radius, log-graded so moderate streets "
     "light up too. The Altstadt retail core glows red — confirms a "
     "dense, contiguous shopping environment."),
    ("proximity_05_transit_stops.jpg", "🚆 Transit stops",
     "Mode-coloured: rail (dark blue), tram (blue), bus (grey). Central "
     "Zürich is tram-saturated; Hauptbahnhof + Central tram hub within "
     "the radius drive the high transit score."),
    ("proximity_06_all_shops.jpg", "🛍️ All shops (generic)",
     "Raw retail footprint — every OSM shop=* (purple) regardless of "
     "brand classification, ~800+ in radius. Shows total commercial "
     "intensity before brand-aware scoring."),
    ("proximity_07_cannibalization.jpg", "⚠️ Cannibalization overlap",
     "Catchment circles + connector lines to existing On stores, "
     "colour-coded by overlap. Niederdorfstrasse 21 is ~200 m from the "
     "Limmatquai flagship → HIGH overlap (red), the key risk here."),
    ("proximity_08_poi_markers.jpg", "🔵 Classified POI markers",
     "The 6 brand-aware categories isolated — sport (blue), symbiose "
     "(green), premium (gold), competitor (red), negative env. (orange), "
     "partner (grey). The core proximity-intelligence layer."),
    ("proximity_09_competitor_circles.jpg", "⚔️ Competitor catchments",
     "Each red 200 m circle = a detected or curated competitor "
     "(Foot Locker, Snipes, Titolo, Adidas Bahnhofstrasse). Explicit "
     "catchment rings rather than a blurred heatmap."),
    ("proximity_10_walk_amenity_heatmap.jpg", "🟢 Walk-amenity density",
     "Food / grocery / services / education / leisure / retail amenities "
     "as a green density heatmap — the underlying data behind the "
     "Walk Score (Niederdorf ≈ 71/100 after May 2026 tuning)."),
    ("proximity_11_sport_heatmap.jpg", "🏋️ Sport-facility density",
     "Gyms, studios, courts, pools as a blue density heatmap. Maps where "
     "On's performance customer actually trains — a core symbiose signal "
     "for the brand."),
    ("proximity_12_hotels.jpg", "🏨 Curated hotels",
     "5★ (green) / 4★ (orange) dots from the 19-hotel curated list "
     "(OSM under-tags stars). Premium-environment proxy feeding the "
     "PEI sub-score."),
]


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, left, top, w, h, fill_rgb, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if not line:
        shape.line.fill.background()
    return shape


def _text(slide, left, top, w, h, text, *, size=14, bold=False,
           color=PRIMARY, align="left"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": 1, "center": 2, "right": 3}.get(align, 1)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _header_band(slide, brand: dict, page_num: int, total: int):
    _rect(slide, 0, 0, SLIDE_W, Inches(0.5), PRIMARY)
    _text(slide, Inches(0.4), Inches(0.08), Inches(9), Inches(0.36),
           f"{brand['brand']}  ·  Retail Site Intelligence  ·  Product Overview",
           size=11, color=WHITE)
    _text(slide, Inches(11.0), Inches(0.08), Inches(2), Inches(0.36),
           f"{page_num} / {total}", size=11, color=WHITE, align="right")


def _slide_title(slide, title: str, subtitle: str = ""):
    _text(slide, Inches(0.4), Inches(0.62), Inches(12.5), Inches(0.5),
           title, size=22, bold=True, color=PRIMARY)
    if subtitle:
        _text(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.4),
               subtitle, size=12, color=SUBTLE)
    _rect(slide, Inches(0.4), Inches(1.4), Inches(0.6), Inches(0.05), ACCENT)


def _embed_screenshot(slide, image_path: Path, *,
                       top_inches: float = 1.6, max_height: float = 5.4) -> None:
    """Add the screenshot to the slide, scaled to fit the available area."""
    if not image_path.exists():
        _text(slide, Inches(0.4), Inches(top_inches),
               Inches(12.5), Inches(0.5),
               f"⚠ Screenshot missing: {image_path.name}",
               size=12, color=PRIMARY)
        return

    with PILImage.open(image_path) as img:
        w_px, h_px = img.size

    # Convert px → inches at 96 DPI (PPTX assumption)
    aspect = h_px / w_px

    # Try fitting by width first
    max_w_in = 12.5
    width_in = max_w_in
    height_in = width_in * aspect

    # If that overflows the area, shrink by height
    if height_in > max_height:
        height_in = max_height
        width_in = height_in / aspect

    # Centre horizontally
    left = Inches((SLIDE_W.inches - width_in) / 2)
    slide.shapes.add_picture(
        str(image_path), left, Inches(top_inches),
        width=Inches(width_in), height=Inches(height_in),
    )


def _footer(slide, caption: str):
    _rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.45), LIGHT)
    _text(slide, Inches(0.4), Inches(7.13), Inches(12.5), Inches(0.32),
           caption, size=9, color=SUBTLE)


# ---------------------------------------------------------------------------
# Specific slides
# ---------------------------------------------------------------------------

def _cover(prs, brand: dict, total: int):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.6), PRIMARY)
    _text(slide, Inches(0.4), Inches(0.12), Inches(8), Inches(0.4),
           f"{brand['brand']}  ·  Retail Site Intelligence",
           size=14, color=WHITE)
    _text(slide, Inches(11.0), Inches(0.12), Inches(2), Inches(0.4),
           datetime.now().strftime("%Y-%m-%d"),
           size=12, color=WHITE, align="right")

    _text(slide, Inches(0.5), Inches(1.4), Inches(12), Inches(0.8),
           "Product Overview — Live Demo Tour",
           size=38, bold=True, color=PRIMARY)
    _text(slide, Inches(0.5), Inches(2.3), Inches(12), Inches(0.6),
           "Site selection · pipeline tracking · portfolio management",
           size=18, color=SUBTLE)

    _rect(slide, Inches(0.5), Inches(3.1), Inches(0.6), Inches(0.07), ACCENT)
    _text(slide, Inches(0.5), Inches(3.2), Inches(12), Inches(0.45),
           "Walk-through: Niederdorfstrasse 21, 8001 Zürich",
           size=20, bold=True, color=PRIMARY)
    _text(slide, Inches(0.5), Inches(3.7), Inches(12), Inches(0.4),
           "Altstadt pedestrian zone · 200 m east of On's Limmatquai flagship · "
           "cannibalization + flood + multi-brand competition all firing",
           size=12, color=SUBTLE)

    # Methodology badges
    badges = [
        ("OSM / OSMnx",      "POI fetch + classification"),
        ("BFS STATPOP",      "Demographics (24 ZH munis)"),
        ("BAFU",             "Flood-risk overlay + override"),
        ("ESI · CCRS",       "Sustainability framework"),
        ("Bechtiger 2024",   "Empirical rent coefficients"),
        ("Glass Box",        "Every dimension inspectable"),
    ]
    y = Inches(4.7)
    col_w = Inches(2.05)
    for i, (label, body) in enumerate(badges):
        x = Inches(0.5) + col_w * i + Inches(0.05 * i)
        _rect(slide, x, y, col_w, Inches(1.6), LIGHT)
        _rect(slide, x, y, Inches(0.06), Inches(1.6), ACCENT)
        _text(slide, x + Inches(0.15), y + Inches(0.12),
               col_w - Inches(0.3), Inches(0.45),
               label, size=11, bold=True, color=PRIMARY)
        _text(slide, x + Inches(0.15), y + Inches(0.6),
               col_w - Inches(0.3), Inches(0.95),
               body, size=9.5, color=SUBTLE)

    _rect(slide, 0, Inches(7.0), SLIDE_W, Inches(0.5), PRIMARY)
    _text(slide, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.32),
           "Generated by the Retail Site Intelligence demo · "
           f"{datetime.now().strftime('%Y-%m-%d %H:%M')}",
           size=9, color=WHITE)


def _shot_slide(prs, brand: dict, page_num: int, total: int,
                 title: str, subtitle: str, screenshot: str, caption: str):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _header_band(slide, brand, page_num, total)
    _slide_title(slide, title, subtitle)
    _embed_screenshot(slide, SCREENSHOTS / screenshot,
                       top_inches=1.6, max_height=5.4)
    _footer(slide, caption)


def _layer_grid_slide(prs, brand: dict, page_num: int, total: int,
                       group_title: str, group_index: int,
                       items: list[tuple[str, str, str]]):
    """One slide with a 2×2 grid of layer screenshots, each with a
    bold layer name + description in a caption box beside the image."""
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _header_band(slide, brand, page_num, total)
    _slide_title(
        slide,
        f"Proximity layers · Niederdorfstrasse 21  ({group_index}/3)",
        group_title,
    )

    # 2×2 grid. Each quadrant: image (left) + caption (right).
    # Image box 3.5×2.70 in (matches 1560×1202 aspect ≈ 1.30).
    img_w = 3.5
    img_h = img_w * (1202 / 1560)   # ≈ 2.70
    positions = [
        (0.35, 1.65),    # top-left
        (6.95, 1.65),    # top-right
        (0.35, 4.45),    # bottom-left
        (6.95, 4.45),    # bottom-right
    ]
    for (img_x, img_y), (fname, name, desc) in zip(positions, items):
        path = LAYER_SHOTS / fname
        if path.exists():
            slide.shapes.add_picture(
                str(path), Inches(img_x), Inches(img_y),
                width=Inches(img_w), height=Inches(img_h),
            )
            # thin border around the image
            border = _rect(slide, Inches(img_x), Inches(img_y),
                            Inches(img_w), Inches(img_h), WHITE)
            border.fill.background()
            border.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            border.line.width = Pt(0.75)
        else:
            _text(slide, Inches(img_x), Inches(img_y), Inches(img_w),
                   Inches(img_h), f"⚠ missing: {fname}", size=10, color=PRIMARY)

        cap_x = img_x + img_w + 0.15
        cap_w = 6.6 - img_w - 0.2     # remaining half-slide width
        # Layer name (bold accent-coloured)
        tb = slide.shapes.add_textbox(
            Inches(cap_x), Inches(img_y), Inches(cap_w), Inches(img_h))
        tf = tb.text_frame
        tf.word_wrap = True
        p_name = tf.paragraphs[0]
        r_name = p_name.add_run()
        r_name.text = name
        r_name.font.size = Pt(12)
        r_name.font.bold = True
        r_name.font.color.rgb = PRIMARY
        p_desc = tf.add_paragraph()
        r_desc = p_desc.add_run()
        r_desc.text = desc
        r_desc.font.size = Pt(9.5)
        r_desc.font.color.rgb = SUBTLE


def _methodology(prs, brand: dict, page_num: int, total: int):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _header_band(slide, brand, page_num, total)
    _slide_title(slide, "Methodology · Data sources · Honest trade-offs", "")

    items = [
        ("OpenStreetMap (OSMnx)",
         "Single query covers sport venues, retail brands, hotels, transit "
         "stops; classification is brand-aware (Lululemon → symbiose, Nike → "
         "competitor) with the 12-step priority cascade in SCORING.md §1."),
        ("Curated competitor overrides",
         "OSM under-tags mono-brand competitor stores in CH — analyst-verified "
         "list catches Adidas Bahnhofstrasse 56, Foot Locker × 3, Snipes × 3, "
         "Titolo. Nike + Puma Sihlcity exits documented for narrative."),
        ("BFS STATPOP + Kaufkraft",
         "Population, growth, age structure, purchasing-power index seeded "
         "for 24 ZH-relevant municipalities — extensible CSV, no auth."),
        ("BAFU surface-runoff hazard",
         "WMS overlay on the map + analyst override for the 4 real On CH "
         "sites."),
        ("ESI (CCRS/UZH) + Bechtiger 2024",
         "Sustainability factors with empirical coefficients: accessibility "
         "(p<.001), daylight (p=.021), ÖV class (n.s. for rent). n = 288 "
         "Swiss income properties · MAS Real Estate UZH."),
        ("Leasehold Pro-Forma",
         "5-yr cashflow for a leased flagship (NOT a property purchase). "
         "Industry KPIs: Occupancy Cost Ratio · Cash-on-Cash · Break-even "
         "sales · Total rent over lease · rent escalation slider."),
        ("Glass Box scoring",
         "Every dimension normalised 0-100 before weighting. Flood penalty "
         "post-aggregation. Sustainability Δ additive outside the 100-pt base."),
        ("Honest trade-offs",
         "Pandana → amenity-richness · BAFU identify API doesn't exist for "
         "the runoff layer → override + WMS · Strava global heatmap "
         "requires auth tokens that expire every ~2 weeks."),
    ]
    y = Inches(1.7)
    for title_item, body in items:
        _rect(slide, Inches(0.4), y, Inches(0.06), Inches(0.6), ACCENT)
        _text(slide, Inches(0.65), y, Inches(4.0), Inches(0.32),
               title_item, size=11, bold=True, color=PRIMARY)
        _text(slide, Inches(0.65), y + Inches(0.3),
               Inches(12.3), Inches(0.4),
               body, size=9, color=SUBTLE)
        y += Inches(0.62)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

SLIDE_SPECS = [
    # title, subtitle, screenshot file, footer caption
    ("Step 1 · Location · BAFU flood-risk + premium-hotel overlays",
     "Geocoded site · cannibalization preview · 🌊 BAFU flood map · "
     "🏨 curated 5★/4★ hotels",
     "02_step1_location.png",
     "Niederdorfstrasse 21 is 200 m east of the existing On Limmatquai "
     "flagship · still inside the Sihl/Limmat HQ100 flood perimeter · "
     "On-site Snipes Niederdorf already in the curated competitor list."),

    ("Step 2 · Market overview · BFS demographics + premium-hotel signal",
     "Population · 5-yr growth · Kaufkraft · age 18-45 · curated hotel count "
     "in radius",
     "03_step2_market.png",
     "Zürich Kaufkraft index 118 (CH avg 100) · 38.5% age 18-45 share · "
     "16 500 hotel beds · curated 5★/4★ hotel count side-by-side with OSM."),

    ("Step 3 · Proximity intelligence · 200+ classified POIs",
     "Sport · Symbiose · Premium · Competitor · Negative · Partner · "
     "+ curated overrides + analytical heatmaps",
     "04_step3_proximity.png",
     "12-step name-then-tag classification cascade. Always-visible curated "
     "competitor table catches Adidas, Foot Locker × 3, Snipes × 3, Titolo "
     "— the real Zürich-city competition that OSM under-tags."),

    ("Step 4 · Footfall · Amenity-richness Walk Score + transit",
     "Tuned log-saturation curve · 6 walk-amenity categories as map layers · "
     "Niederdorfstr. Walk ~71/100, Transit 100/100",
     "05_step4_footfall.png",
     "Pandana-free implementation, May 2026 saturation tuned to "
     "13·ln(1+s) so central Altstadt no longer trivially saturates "
     "at 100. Realistic differentiation: Bahnhofstr 72, Niederdorf 71, "
     "Hardturm 51, Wipkingen 30."),

    ("Step 5 · Site Score · Glass Box methodology",
     "7-dimension normalised 0-100 score · BAFU flood penalty + Bechtiger "
     "ESI Δ applied post-aggregation",
     "06_step5_score.png",
     "Glass-Box breakdown shows what raises / lowers the score. Flood "
     "applied post-aggregation; sustainability Δ additive outside the "
     "100-pt base. Site Overview Map at the bottom is the analyst's "
     "synthesis view."),

    ("Step 6 · Leasehold Pro-Forma · 5-yr cashflow + retail KPIs",
     "NPV · IRR · Payback · Occupancy Cost Ratio · Cash-on-Cash · "
     "Break-even sales · Bear/Base/Bull",
     "07_step6_proforma.png",
     "Leasehold model (no property purchase). Industry-standard KPIs "
     "for retail flagships: OCR target <10%, Cash-on-Cash Y1 indicates "
     "fit-out payback speed, rent escalation slider models Swiss "
     "commercial-lease LIK indexation."),

    ("Step 7 · Export · PDF Site Report + Deal Memo + PPT",
     "1-pager · 5-page deal memo · 6-slide product overview · all assembled "
     "from real data",
     "08_step7_export.png",
     "ReportLab PDFs with embedded OSM-basemap maps, charts, top-POI tables. "
     "python-pptx product overview deck. All assembled from the current "
     "session's analysis state — no static placeholders."),
]


def build_screenshot_deck(brand: dict) -> bytes:
    """Build a slide deck with embedded Playwright screenshots.

    Structure: cover · 7 step slides · 3 proximity-layer-detail slides
    (4 layers each) · methodology.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 3 layer-detail slides, 4 layers each (12 overlays)
    layer_groups = [
        ("Environment & access · flood · running · pedestrian · shop density",
         PROXIMITY_LAYERS[0:4]),
        ("Mobility & retail footprint · transit · all shops · cannibalization · POI split",
         PROXIMITY_LAYERS[4:8]),
        ("Brand-fit signals · competitors · walk amenities · sport density · hotels",
         PROXIMITY_LAYERS[8:12]),
    ]

    total = 1 + len(SLIDE_SPECS) + len(layer_groups) + 1   # cover+steps+layers+methodology
    _cover(prs, brand, total)

    page = 2
    for (title, subtitle, shot, caption) in SLIDE_SPECS:
        _shot_slide(prs, brand, page, total, title, subtitle, shot, caption)
        page += 1

    for gi, (group_title, items) in enumerate(layer_groups, 1):
        _layer_grid_slide(prs, brand, page, total, group_title, gi, items)
        page += 1

    _methodology(prs, brand, page, total)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
